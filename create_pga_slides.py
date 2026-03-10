from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_NAVY = RGBColor(26, 26, 46)
GOLD = RGBColor(212, 175, 55)
GREEN = RGBColor(26, 71, 42)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(248, 249, 250)

def add_background(slide, color):
    """Add solid background to slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_slide(prs, title, subtitle, header_text=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, DARK_NAVY)
    
    # Gold accent bar at top
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLD
    top_bar.line.fill.background()
    
    # PGA header if provided
    if header_text:
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.5))
        header_frame = header_box.text_frame
        header_frame.text = header_text
        header_frame.paragraphs[0].font.size = Pt(14)
        header_frame.paragraphs[0].font.color.rgb = GOLD
        header_frame.paragraphs[0].font.bold = True
        header_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.8), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.8), Inches(1))
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        tf2.text = subtitle
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(180, 180, 200)
        p2.alignment = PP_ALIGN.CENTER
    
    # Bottom gold accent
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = GOLD
    bottom_bar.line.fill.background()
    
    return slide

def add_content_slide(prs, title, bullets, header=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, WHITE)
    
    # Green header bar
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = GREEN
    header_bar.line.fill.background()
    
    # Title on header
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Optional header text
    if header:
        header_txt = slide.shapes.add_textbox(Inches(0.6), Inches(0.05), Inches(12), Inches(0.3))
        htf = header_txt.text_frame
        htf.text = header
        hp = htf.paragraphs[0]
        hp.font.size = Pt(11)
        hp.font.color.rgb = GOLD
        hp.font.bold = True
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(5.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(50, 50, 60)
        p.space_after = Pt(16)
        p.level = 0
    
    return slide

def add_split_slide(prs, title, left_title, left_items, right_title, right_items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, WHITE)
    
    # Header
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = DARK_NAVY
    header_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.6))
    ltf = left_title_box.text_frame
    ltf.text = left_title
    lp = ltf.paragraphs[0]
    lp.font.size = Pt(24)
    lp.font.bold = True
    lp.font.color.rgb = RGBColor(150, 150, 160)
    lp.font.italic = True
    
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.5))
    ltf2 = left_box.text_frame
    ltf2.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = ltf2.paragraphs[0]
        else:
            p = ltf2.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(100, 100, 110)
        p.space_after = Pt(12)
    
    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.6), Inches(0.02), Inches(5.2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 220, 230)
    line.line.fill.background()
    
    # Right column
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(0.6))
    rtf = right_title_box.text_frame
    rtf.text = right_title
    rp = rtf.paragraphs[0]
    rp.font.size = Pt(24)
    rp.font.bold = True
    rp.font.color.rgb = GREEN
    
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.3), Inches(5.8), Inches(4.5))
    rtf2 = right_box.text_frame
    rtf2.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = rtf2.paragraphs[0]
        else:
            p = rtf2.add_paragraph()
        p.text = f"✓ {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(26, 71, 42)
        p.font.bold = True
        p.space_after = Pt(14)
    
    return slide

def add_numbered_slide(prs, title, number, subtitle, description):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, RGBColor(248, 249, 250))
    
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), Inches(0.6), Inches(1.2), Inches(1.2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(1.2), Inches(0.7))
    ntf = num_box.text_frame
    ntf.text = str(number)
    np = ntf.paragraphs[0]
    np.font.size = Pt(40)
    np.font.bold = True
    np.font.color.rgb = GOLD
    np.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(2.3), Inches(0.7), Inches(10.5), Inches(0.9))
    tf = title_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    dtf = desc_box.text_frame
    dtf.word_wrap = True
    dtf.text = description
    dp = dtf.paragraphs[0]
    dp.font.size = Pt(26)
    dp.font.color.rgb = RGBColor(60, 60, 70)
    dp.line_spacing = 1.3
    
    # Gold accent bar bottom
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.2), prs.slide_width, Inches(0.3)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = GOLD
    bottom_bar.line.fill.background()
    
    return slide

def add_closing_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, GREEN)
    
    # Gold accent bars
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLD
    top_bar.line.fill.background()
    
    # Main message
    msg_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.8), Inches(1.5))
    tf = msg_box.text_frame
    tf.word_wrap = True
    tf.text = "One Leadership Shift"
    p = tf.paragraphs[0]
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Arrow
    arrow_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.5), Inches(11.8), Inches(0.8))
    atf = arrow_box.text_frame
    atf.text = "↓"
    ap = atf.paragraphs[0]
    ap.font.size = Pt(60)
    ap.font.color.rgb = GOLD
    ap.alignment = PP_ALIGN.CENTER
    
    # Result
    result_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.2), Inches(11.8), Inches(2.5))
    rtf = result_box.text_frame
    rtf.word_wrap = True
    rtf.text = "Consistent Member Experience\nStronger Loyalty & Retention"
    rp = rtf.paragraphs[0]
    rp.font.size = Pt(36)
    rp.font.bold = True
    rp.font.color.rgb = GOLD
    rp.alignment = PP_ALIGN.CENTER
    
    # Bottom bar
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.3), prs.slide_width, Inches(0.2)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = GOLD
    bottom_bar.line.fill.background()
    
    return slide

# ============ BUILD SLIDES ============

# Slide 1: Title
add_title_slide(
    prs,
    "The Hospitality\nLeadership Mechanism",
    "How Top Clubs Elevate Member Experience in 2026",
    "PGA OF ONTARIO • THE LEADERSHIP ADVANTAGE BEHIND EXCEPTIONAL MEMBER EXPERIENCES"
)

# Slide 2: This is NOT...
add_content_slide(
    prs,
    "What This Is",
    [
        "This is NOT a service training",
        "It IS a leadership mechanism for Pros and key team members",
        "Equips leaders to elevate member experience every day — even with young or seasonal staff",
        "Participants walk away with a repeatable system they can deploy immediately",
        "Creates a culture of hospitality that sustains itself beyond any single season"
    ],
    "BREAKOUT WORKSHOP"
)

# Slide 3: The 5 Accelerators Overview
add_content_slide(
    prs,
    "The 5 Hospitality Leadership Accelerators",
    [
        "The core leadership behaviours that create self-sustaining hospitality culture",
        "These become the mechanism leaders use to raise standards without micromanaging",
        "Each accelerator builds on the others to create a complete leadership system",
        "Designed for immediate application in real club environments"
    ],
    "THE FRAMEWORK"
)

# Slides 4-8: Individual Accelerators
accelerators = [
    ("Leadership Certainty", "Presence, confidence, and composure that sets the tone for the entire team and member experience. When leaders show up with certainty, everyone else rises to match that energy."),
    ("Impactful Communication", "Clarity, tone, influence, and listening — the skills that ensure messages land and motivate action. Great leaders communicate in ways that move people to action, not just understanding."),
    ("Human Influence & Relationship Building", "Gaining cooperation and shifting mindset without friction — the art of moving people willingly. When people feel connected, they contribute more than what's required."),
    ("Leadership in Action", "Moving people to places they wouldn't go on their own — inspiring discretionary effort and ownership. True leadership is measured by what people do when no one is watching."),
    ("Energy & Emotional State Management", "Controlling attitude, elevating enthusiasm, and inspiring others — regardless of circumstances. Leaders set the emotional temperature of every interaction.")
]

for i, (title, desc) in enumerate(accelerators, 1):
    add_numbered_slide(prs, "Accelerator", i, title, desc)

# Slide 9: Bridge Leadership Technique
add_content_slide(
    prs,
    "The Bridge Leadership Technique",
    [
        "A simple, repeatable coaching tool leaders can use immediately",
        "Lifts performance and behaviour without confrontation or defensiveness",
        "Applied to real club scenarios in the workshop:",
        "  → First Tee First Impressions — Setting the tone in the first 3 minutes",
        "  → Fast Recovery After a Service Breakdown — Resetting without defensiveness",
        "  → Coaching a Disengaged Seasonal Employee — Improving ownership without conflict",
        "A 'use tomorrow' skill that sticks"
    ],
    "PRACTICAL TOOLS"
)

# Slide 10: Comparison
add_split_slide(
    prs,
    "Why This Works",
    "Traditional Approach",
    [
        "• Train staff on service behaviours",
        "• Correct issues after they happen",
        "• Culture resets every season",
        "• Reactive problem-solving",
        "• Dependent on constant oversight"
    ],
    "The Mechanism",
    [
        "Develop leaders who create a culture people rise to",
        "Prevent breakdowns through proactive emotional leadership",
        "Culture becomes durable, repeatable, and staff-led",
        "Proactive culture-building",
        "Self-sustaining system"
    ]
)

# Slide 11: The Shift
add_split_slide(
    prs,
    "The Leadership Shift",
    "From",
    [
        "• Training events",
        "• Managing tasks",
        "• Fixing problems",
        "• Supervising staff",
        "• Seasonal resets"
    ],
    "To",
    [
        "Leadership systems",
        "Developing people",
        "Building culture",
        "Coaching leaders",
        "Sustainable excellence"
    ]
)

# Slide 12: About Troy
add_content_slide(
    prs,
    "About Your Facilitator",
    [
        "Troy Treleaven, CEO — Dale Carnegie Training Ontario & Maritimes",
        "Decades of experience developing leaders across industries",
        "Specializes in sustainable leadership systems that drive measurable results",
        "Work with golf clubs, resorts, and hospitality organizations",
        "Focus: Intersection of leadership behaviour and customer experience",
        "Creating cultures where exceptional service is the natural outcome"
    ],
    "TROY TRELEAVEN"
)

# Slide 13: Event Details
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
add_background(slide, WHITE)

header_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3)
)
header_bar.fill.solid()
header_bar.fill.fore_color.rgb = DARK_NAVY
header_bar.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "Event Details"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

details = [
    ("Date:", "Monday, February 23, 2026"),
    ("Venue:", "TPC Toronto at Osprey Valley"),
    ("Address:", "19131 Main St, Alton, ON L7K 1R1"),
    ("Entrance:", "Main doors — PGA team on-site from 7:30 AM"),
    ("Arrival:", "Registrants may arrive as early as 8:30 AM"),
    ("", ""),
    ("Session 1:", "9:30 AM – 10:30 AM"),
    ("Session 2:", "10:40 AM – 11:40 AM"),
    ("Session 3:", "11:50 AM – 12:50 PM")
]

y_pos = 1.7
for label, value in details:
    if label:
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(2.5), Inches(0.5))
        ltf = lbl_box.text_frame
        ltf.text = label
        lp = ltf.paragraphs[0]
        lp.font.size = Pt(18)
        lp.font.bold = True
        lp.font.color.rgb = GREEN
        
        # Value
        val_box = slide.shapes.add_textbox(Inches(3.3), Inches(y_pos), Inches(9), Inches(0.5))
        vtf = val_box.text_frame
        vtf.text = value
        vp = vtf.paragraphs[0]
        vp.font.size = Pt(18)
        vp.font.color.rgb = RGBColor(50, 50, 60)
    y_pos += 0.55

# Slide 14: Closing
add_closing_slide(prs)

# Save
output_path = "/data/.openclaw/workspace/PGA_Ontario_Hospitality_Leadership_Workshop.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
